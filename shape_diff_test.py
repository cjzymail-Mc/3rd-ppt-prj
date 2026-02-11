import os
import json
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def get_shape_geometry_key(shape):
    return (
        shape.shape_type,
        int(round(shape.left, -4)),
        int(round(shape.top, -4)),
        int(round(shape.width, -4)),
        int(round(shape.height, -4))
    )

def analyze_ppt(path, label=""):
    prs = Presentation(path)
    slide = prs.slides[0]
    shapes = {}
    for shape in slide.shapes:
        key = get_shape_geometry_key(shape)
        if key not in shapes:
            shapes[key] = []
        shapes[key].append(shape)
    return shapes

def get_font_details(shape):
    if not hasattr(shape, "text_frame"):
        return None
    try:
        if shape.text_frame.paragraphs:
            p = shape.text_frame.paragraphs[0]
            if p.runs:
                font = p.runs[0].font
                return {
                    "name": font.name,
                    "size": font.size,
                    "color": str(font.color.rgb) if hasattr(font.color, 'rgb') else "None",
                    "bold": font.bold,
                    "italic": font.italic
                }
    except Exception as e:
        return None
    return None

def update_class_030(font_names):
    class_file = "class_030.py"
    if not os.path.exists(class_file):
        # Create basic file if missing
        with open(class_file, "w") as f:
            f.write("class Text_Box:\n    pass\n\n")

    with open(class_file, "r") as f:
        content = f.read()

    for font in font_names:
        if not font: continue
        # sanitize font name for class name
        safe_name = re.sub(r'[^a-zA-Z0-9]', '_', font)
        class_name = f"Font_{safe_name}"

        if class_name not in content:
            print(f"Adding new font class {class_name} to {class_file}")
            with open(class_file, "a") as f:
                f.write(f"\nclass {class_name}:\n    NAME = \"{font}\"\n")

def main():
    blank_path = "debug/篮球试穿问卷ppt模板 - 0.pptx"
    standard_path = "debug/篮球试穿问卷ppt模板 - 2.pptx"

    if not os.path.exists(blank_path) or not os.path.exists(standard_path):
        print("Error: Files not found.")
        return

    blank_shapes_map = analyze_ppt(blank_path, "Blank")
    standard_shapes_map = analyze_ppt(standard_path, "Standard")

    new_elements = []

    for key, shapes_std in standard_shapes_map.items():
        if key not in blank_shapes_map:
            for s in shapes_std:
                new_elements.append(s)
        else:
            count_std = len(shapes_std)
            count_blank = len(blank_shapes_map[key])
            if count_std > count_blank:
                diff = count_std - count_blank
                for i in range(diff):
                    new_elements.append(shapes_std[count_blank + i])

    print(f"Found {len(new_elements)} new elements.")

    elements_data = []
    found_fonts = set()

    for i, shape in enumerate(new_elements):
        font_info = get_font_details(shape)
        content = shape.text if hasattr(shape, "text") else ""

        if font_info and font_info["name"]:
            found_fonts.add(font_info["name"])

        element = {
            "index": i + 1,
            "id": shape.shape_id,
            "name": shape.name,
            "type_id": int(shape.shape_type),
            "top": shape.top,
            "left": shape.left,
            "width": shape.width,
            "height": shape.height,
            "content": content,
            "font": font_info
        }
        elements_data.append(element)

    # Update class_030.py
    update_class_030(found_fonts)

    with open("shapes.json", "w", encoding="utf-8") as f:
        json.dump(elements_data, f, indent=2, ensure_ascii=False)

    with open("shape-detail.md", "w", encoding="utf-8") as f:
        f.write("# Shape Details (New Elements)\n\n")
        for d in elements_data:
            f.write(f"## Element {d['index']}\n")
            f.write(f"- **ID**: {d['id']}\n")
            f.write(f"- **Name**: {d['name']}\n")
            f.write(f"- **Type**: {d['type_id']}\n")
            f.write(f"- **Position**: Top={d['top']}, Left={d['left']}\n")
            f.write(f"- **Size**: Height={d['height']}, Width={d['width']}\n")
            f.write(f"- **Content**: {d['content']}\n")
            if d['font']:
                f.write(f"- **Font**: {d['font']}\n")
            f.write("\n")

if __name__ == "__main__":
    main()
