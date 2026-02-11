import json
import copy
import pandas as pd
from pptx import Presentation
from pptx.oxml import parse_xml

def get_excel_data():
    df = pd.read_excel("debug/问卷模板 - 2.xlsx")
    # Take first row
    row = df.iloc[0]
    return {
        "title": str(row.get("试穿鞋款名称", "Unknown")).upper().replace("PRO", " PRO"),
        "补充说明": str(row.get("补充说明", "")),
        "你认为这双鞋更加适合什么打法的球员穿": str(row.get("你认为这双鞋更加适合什么打法的球员穿", "")),
        "抓地性（Traction）": float(row.get("抓地性（Traction）", 0)),
        "包裹性（Lockdwon）": float(row.get("包裹性（Lockdwon）", 0)),
        "重量（Weight）": float(row.get("重量（Weight）", 0)),
        "缓震性（Cushioning）": float(row.get("缓震性（Cushioning）", 0)),
        "防侧翻性（Lateral Stability）": float(row.get("防侧翻性（Lateral Stability）", 0)),
        "耐久性（Durability）": float(row.get("耐久性（Durability）", 0))
    }

def add_title(shapes_map, data):
    if 2 in shapes_map and hasattr(shapes_map[2], "text"):
        shapes_map[2].text = data.get("title", "")

def add_basic_info(shapes_map, data):
    if 6 in shapes_map and hasattr(shapes_map[6], "text"): shapes_map[6].text = "第一轮"
    if 8 in shapes_map and hasattr(shapes_map[8], "text"): shapes_map[8].text = "| 竞技&训练"
    if 10 in shapes_map and hasattr(shapes_map[10], "text"): shapes_map[10].text = "5人"
    if 14 in shapes_map and hasattr(shapes_map[14], "text"): shapes_map[14].text = "2人"
    if 16 in shapes_map and hasattr(shapes_map[16], "text"): shapes_map[16].text = "7人"
    if 18 in shapes_map and hasattr(shapes_map[18], "text"): shapes_map[18].text = "两周 - 2025.10"
    if 19 in shapes_map and hasattr(shapes_map[19], "text"): shapes_map[19].text = ">"

def add_labels(shapes_map, data):
    labels = ["抓地", "包裹", "重量", "缓震", "稳定", "耐久"]
    ids = [49, 50, 51, 52, 53, 54]
    for i, pid in enumerate(ids):
        if pid in shapes_map and hasattr(shapes_map[pid], "text"):
            shapes_map[pid].text = labels[i]

def add_scores(shapes_map, data):
    scores = [
        data.get("抓地性（Traction）", 0),
        data.get("包裹性（Lockdwon）", 0),
        data.get("重量（Weight）", 0),
        data.get("缓震性（Cushioning）", 0),
        data.get("防侧翻性（Lateral Stability）", 0),
        data.get("耐久性（Durability）", 0)
    ]
    ids = [56, 58, 60, 62, 64, 65]
    for i, pid in enumerate(ids):
        if pid in shapes_map and hasattr(shapes_map[pid], "text"):
            shapes_map[pid].text = f"{float(scores[i]):.1f}"

def calculate_average(data):
    scores = [
        data.get("抓地性（Traction）", 0),
        data.get("包裹性（Lockdwon）", 0),
        data.get("重量（Weight）", 0),
        data.get("缓震性（Cushioning）", 0),
        data.get("防侧翻性（Lateral Stability）", 0),
        data.get("耐久性（Durability）", 0)
    ]
    return sum(scores) / len(scores) if scores else 0

def add_average(shapes_map, data):
    avg = calculate_average(data)
    if 69 in shapes_map and hasattr(shapes_map[69], "text"):
        shapes_map[69].text = f"{avg:.1f}"

def add_grade(shapes_map, data):
    avg = calculate_average(data)
    grade = "C"
    suffix = ""

    if avg >= 9: grade = "S"
    elif avg >= 8: grade = "A"
    elif avg >= 7: grade = "B"
    elif avg >= 6: grade = "C"
    else: grade = "D"

    # Check decimal part for suffix
    dec = avg - int(avg)
    if dec >= 0.7: suffix = "+"
    elif dec <= 0.3: suffix = "-"

    if 78 in shapes_map and hasattr(shapes_map[78], "text"): shapes_map[78].text = grade
    if 79 in shapes_map and hasattr(shapes_map[79], "text"): shapes_map[79].text = suffix

def add_comments(shapes_map, data):
    if 90 in shapes_map and hasattr(shapes_map[90], "text"):
        shapes_map[90].text = data.get("补充说明", "")

def add_player_style(shapes_map, data):
    if 105 in shapes_map and hasattr(shapes_map[105], "text"):
        shapes_map[105].text = data.get("你认为这双鞋更加适合什么打法的球员穿", "")

def add_priority(shapes_map, data):
    pass

def main():
    try:
        data = get_excel_data()

        with open("shapes.json", "r", encoding="utf-8") as f:
            shapes_def = json.load(f)

        target_prs = Presentation("debug/篮球试穿问卷ppt模板 - 0.pptx")
        target_slide = target_prs.slides[0]

        source_prs = Presentation("debug/篮球试穿问卷ppt模板 - 2.pptx")
        source_slide = source_prs.slides[0]

        source_shapes_by_id = {}
        for s in source_slide.shapes:
            source_shapes_by_id[s.shape_id] = s

        new_shapes_map = {}

        for item in shapes_def:
            orig_id = item["id"]
            if orig_id in source_shapes_by_id:
                source_shape = source_shapes_by_id[orig_id]
                new_element = copy.deepcopy(source_shape.element)
                target_slide.shapes._spTree.append(new_element)
                new_shape = target_slide.shapes[-1]
                new_shapes_map[orig_id] = new_shape

        add_title(new_shapes_map, data)
        add_basic_info(new_shapes_map, data)
        add_labels(new_shapes_map, data)
        add_scores(new_shapes_map, data)
        add_average(new_shapes_map, data)
        add_grade(new_shapes_map, data)
        add_comments(new_shapes_map, data)
        add_player_style(new_shapes_map, data)
        add_priority(new_shapes_map, data)

        target_prs.save("codex.pptx")
        print("Successfully generated codex.pptx")

    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()
